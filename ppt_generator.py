from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import difflib


def add_text_to_frame(text_frame, text):
    """Formats text properly: bold for subheadings, italic for sub-subheadings, and avoids overflow."""
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)

    lines = text.split("\n")
    for line in lines:
        if line.strip():
            paragraph = text_frame.add_paragraph()
            run = paragraph.add_run()
            run.text = line.strip()
            run.font.size = Pt(18)

            if line.startswith("## "):
                run.font.bold = True
                run.text = line.replace("## ", "")
            elif line.startswith("### "):
                run.font.italic = True
                run.text = line.replace("### ", "")


def create_presentation(topic, sections, slides_data, images_data, inverted_image_y_positions=None):
    """Creates a PowerPoint presentation with alternating layouts and custom Y-axis for inverted images."""
    presentation = Presentation()

    text_top = Inches(2)
    text_width = Inches(5)
    text_height = Inches(6)

    image_width = Inches(3.5)
    image_height = Inches(3)

    if inverted_image_y_positions is None:
        inverted_image_y_positions = {}

    for i, (section, (slide_text, images)) in enumerate(zip(sections, zip(slides_data, images_data))):
        max_chars_per_slide = 550
        text_parts = [slide_text[i:i + max_chars_per_slide] for i in range(0, len(slide_text), max_chars_per_slide)]

        for part_index, text_part in enumerate(text_parts):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])

            title_box = slide.shapes.title
            if title_box:
                title_box.text = section.upper() if part_index == 0 else f"{section.upper()} "

            if i % 2 == 0:
                text_x = Inches(0.5)
                image_x = Inches(5.5)
                image_y = Inches(2.5)
            else:
                text_x = Inches(5)
                image_x = Inches(1)
                image_y = inverted_image_y_positions.get(i, Inches(2.5))

            content_box = slide.shapes.add_textbox(text_x, text_top, text_width, text_height)
            content_frame = content_box.text_frame
            add_text_to_frame(content_frame, text_part)

            if images:
                image_index = part_index % len(images)
                image_filename = images[image_index]
                slide.shapes.add_picture(image_filename, image_x, image_y, image_width, image_height)

    ppt_path = f"{topic.replace(' ', '_')}.pptx"
    presentation.save(ppt_path)
    return ppt_path
