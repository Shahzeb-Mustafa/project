from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_text_to_frame(text_frame, text):
    """Formats text properly: bold for subheadings, italic for sub-subheadings, and avoids overflow."""
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)

    # Splitting text into sections
    lines = text.split("\n")
    for line in lines:
        if line.strip():
            paragraph = text_frame.add_paragraph()
            run = paragraph.add_run()
            run.text = line.strip()
            run.font.size = Pt(18)

            # Formatting rules
            if line.startswith("## "):  # Subheading
                run.font.bold = True
                run.text = line.replace("## ", "")  # Remove ##
            elif line.startswith("### "):  # Sub-subheading
                run.font.italic = True
                run.text = line.replace("### ", "")  # Remove ###


def create_presentation(topic, sections, slides_data, images_data):
    presentation = Presentation()

    # Define text box dimensions
    text_left = Inches(0.5)
    text_top = Inches(1.5)
    text_width = Inches(4)
    text_height = Inches(6)

    # Define image size
    image_width = Inches(3)
    image_height = Inches(3)

    for i, (section, (slide_text, images)) in enumerate(zip(sections, zip(slides_data, images_data))):
        # Split long text into multiple slides
        max_chars_per_slide = 400
        text_parts = [slide_text[i:i + max_chars_per_slide] for i in range(0, len(slide_text), max_chars_per_slide)]

        for part_index, text_part in enumerate(text_parts):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])

            # Set title (Only for first slide of the section)
            title_box = slide.shapes.title
            if title_box:
                title_box.text = section if part_index == 0 else f"{section} (Cont.)"

            # Add text box (content area)
            content_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            content_frame = content_box.text_frame
            add_text_to_frame(content_frame, text_part)

            # Define image positions (Alternate layouts)
            if i % 2 == 0:
                # Original layout (image top-right)
                image_x = Inches(6.5)
                image_y = Inches(3.0)
            else:
                # Inverted layout (image bottom-left)
                image_x = Inches(0.5)
                image_y = Inches(6.0)

            # Add image (Only on the first slide of the section)
            if part_index == 0:
                slide.shapes.add_picture(images[0], image_x, image_y, image_width, image_height)

    # Save the PowerPoint
    ppt_path = f"{topic.replace(' ', '_')}.pptx"
    presentation.save(ppt_path)
    return ppt_path


