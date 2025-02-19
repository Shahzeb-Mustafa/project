from openai import OpenAI
import os
import time
import requests
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
from PIL import Image
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from gtts import gTTS
import re

def get_topic_details(topic, sections):
    client = OpenAI(
        base_url="https://integrate.api.nvidia.com/v1",
        api_key="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"
    )

    presentation = Presentation()
    lecture_script = "Lecture Script for Topic: " + topic + "\n\n"

    for section in sections:
        prompt = (
            f"Provide a comprehensive, detailed explanation about {topic} focusing on the section: '{section}'. "
            "Include relevant information, examples, and context. Provide content suitable for a presentation slide "
            "and also a detailed lecture script with explanations, examples, and a teaching-friendly format."
        )

        completion = client.chat.completions.create(
            model="meta/llama-3.1-70b-instruct",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2,
            top_p=0.7,
            max_tokens=1024,
            stream=True
        )

        section_details = ""
        for chunk in completion:
            if chunk.choices[0].delta.content is not None:
                section_details += chunk.choices[0].delta.content

        section_details = re.sub(r"\*", "", section_details)

        slide_content, lecture_content = section_details[:500], section_details[500:]

        lecture_script += f"### {section}\n{lecture_content}\n\n"

        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = section

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.text = slide_content

        image_paths = search_and_download_images_bing(f"{topic} {section}", num_images=2)
        for idx, img_path in enumerate(image_paths):
            slide.shapes.add_picture(img_path, Inches(6.5), Inches(1.5 + idx * 3), Inches(3))

    ppt_path = f"{topic.replace(' ', '_')}.pptx"
    presentation.save(ppt_path)
    print(f"Presentation saved as {ppt_path}")

    script_path = f"{topic.replace(' ', '_')}_lecture_script.txt"
    with open(script_path, "w", encoding="utf-8") as f:
        f.write(lecture_script)
    print(f"Lecture script saved as {script_path}")

    audio_path = f"{topic.replace(' ', '_')}_lecture_audio.mp3"
    tts = gTTS(lecture_script)
    tts.save(audio_path)
    print(f"Lecture audio saved as {audio_path}")

def search_and_download_images_bing(query, num_images=2):
    options = webdriver.ChromeOptions()
    options.add_argument('--start-maximized')
    options.add_argument("user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/119.0.0.0 Safari/537.36")

    driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=options)

    search_url = f"https://www.bing.com/images/search?q={query}&qft=+filterui:imagesize-large"
    driver.get(search_url)

    for _ in range(5):
        driver.find_element(By.TAG_NAME, 'body').send_keys(Keys.END)
        time.sleep(2)

    image_elements = driver.find_elements(By.CSS_SELECTOR, "img.mimg")
    image_urls = set()

    for img in image_elements:
        src = img.get_attribute('src') or img.get_attribute('data-src')
        if src and src.startswith('http'):
            image_urls.add(src)
            if len(image_urls) >= num_images:
                break

    driver.quit()

    save_path = os.path.abspath(query.replace(' ', '_'))
    os.makedirs(save_path, exist_ok=True)

    downloaded_images = []

    for i, url in enumerate(image_urls):
        try:
            print(f"Downloading image {i+1} from {url}")
            response = requests.get(url, stream=True, timeout=10)
            response.raise_for_status()
            img = Image.open(BytesIO(response.content))
            img = img.convert('RGB')
            img_path = os.path.join(save_path, f"image_{i+1}.jpg")
            img.save(img_path, quality=95)
            downloaded_images.append(img_path)
            print(f"Saved image_{i+1}.jpg at {save_path}")
        except Exception as e:
            print(f"Failed to save image from {url}: {e}")

    return downloaded_images

def main():
    topic = input("Enter a topic: ")
    sections_input = input("Enter sections separated by commas (e.g., Introduction, History, Applications): ")
    sections = [section.strip() for section in sections_input.split(',')]
    print("\nFetching details, images, and generating presentation with lecture script and audio...\n")
    get_topic_details(topic, sections)
    print(f"\nDetails displayed, images downloaded, presentation, lecture script, and audio created successfully.")

if __name__ == "__main__":
    main()